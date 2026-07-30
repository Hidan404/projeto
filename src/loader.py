import csv
import json
import logging
import os
from typing import List

from langchain_core.documents import Document

logger = logging.getLogger(__name__)


def carregar_documento(caminho_arquivo: str) -> List[Document]:
    ext = os.path.splitext(caminho_arquivo)[1].lower()
    mapa_ext = {
        ".pdf": _carregar_pdf,
        ".txt": _carregar_texto,
        ".md": _carregar_texto,
        ".csv": _carregar_csv,
        ".json": _carregar_json,
        ".docx": _carregar_docx,
        ".pptx": _carregar_pptx,
        ".html": _carregar_html,
        ".htm": _carregar_html,
        ".xlsx": _carregar_xlsx,
        ".xls": _carregar_xlsx,
    }
    func_carregador = mapa_ext.get(ext)
    if not func_carregador:
        raise ValueError(f"Formato nao suportado: {ext}")
    return func_carregador(caminho_arquivo)


def carregar_documentos_do_diretorio(diretorio: str) -> List[Document]:
    todos_docs = []
    for raiz, _, arquivos in os.walk(diretorio):
        for nome_arquivo in arquivos:
            caminho = os.path.join(raiz, nome_arquivo)
            try:
                docs = carregar_documento(caminho)
                todos_docs.extend(docs)
                logger.info("%s (%d paginas/blocos)", nome_arquivo, len(docs))
            except Exception as e:
                logger.warning("%s: %s", nome_arquivo, e)
    return todos_docs


def _carregar_pdf(caminho: str) -> List[Document]:
    from pypdf import PdfReader
    docs = []
    leitor = PdfReader(caminho)
    for i, pagina in enumerate(leitor.pages):
        texto = pagina.extract_text()
        if texto.strip():
            docs.append(Document(
                page_content=texto,
                metadata={"source": caminho, "pagina": i + 1, "total_paginas": len(leitor.pages)},
            ))
    return docs


def _carregar_texto(caminho: str) -> List[Document]:
    with open(caminho, "r", encoding="utf-8") as f:
        texto = f.read()
    return [Document(page_content=texto, metadata={"source": caminho})]


def _carregar_csv(caminho: str) -> List[Document]:
    docs = []
    with open(caminho, "r", encoding="utf-8") as f:
        leitor = csv.DictReader(f)
        for i, linha in enumerate(leitor):
            conteudo = "\n".join(f"{k}: {v}" for k, v in linha.items())
            docs.append(Document(
                page_content=conteudo,
                metadata={"source": caminho, "linha": i + 1},
            ))
    return docs


def _carregar_json(caminho: str) -> List[Document]:
    with open(caminho, "r", encoding="utf-8") as f:
        try:
            dados = json.load(f)
        except json.JSONDecodeError as e:
            raise ValueError(f"Arquivo JSON mal formado: {e}") from e
    texto = json.dumps(dados, ensure_ascii=False, indent=2)
    return [Document(page_content=texto, metadata={"source": caminho})]


def _carregar_docx(caminho: str) -> List[Document]:
    from docx import Document as DocxDocument
    try:
        doc = DocxDocument(caminho)
    except Exception as e:
        raise ValueError(f"Arquivo DOCX corrompido ou invalido: {e}") from e
    paragrafos = [p.text for p in doc.paragraphs if p.text.strip()]
    texto = "\n".join(paragrafos)
    return [Document(page_content=texto, metadata={"source": caminho})]


def _carregar_pptx(caminho: str) -> List[Document]:
    from pptx import Presentation
    try:
        apresentacao = Presentation(caminho)
    except Exception as e:
        raise ValueError(f"Arquivo PPTX corrompido ou invalido: {e}") from e
    docs = []
    for i, slide in enumerate(apresentacao.slides):
        textos = []
        for forma in slide.shapes:
            if forma.has_text_frame:
                for paragrafo in forma.text_frame.paragraphs:
                    if paragrafo.text.strip():
                        textos.append(paragrafo.text)
        if textos:
            docs.append(Document(
                page_content="\n".join(textos),
                metadata={"source": caminho, "slide": i + 1},
            ))
    return docs


def _carregar_html(caminho: str) -> List[Document]:
    from bs4 import BeautifulSoup
    with open(caminho, "r", encoding="utf-8") as f:
        try:
            soup = BeautifulSoup(f.read(), "lxml")
        except Exception as e:
            raise ValueError(f"Arquivo HTML mal formado: {e}") from e
    texto = soup.get_text(separator="\n", strip=True)
    return [Document(page_content=texto, metadata={"source": caminho})]


def _carregar_xlsx(caminho: str) -> List[Document]:
    import pandas as pd
    try:
        planilha = pd.ExcelFile(caminho)
    except Exception as e:
        raise ValueError(f"Arquivo XLSX corrompido ou invalido: {e}") from e
    docs = []
    for nome_aba in planilha.sheet_names:
        df = pd.read_excel(caminho, sheet_name=nome_aba)
        conteudo = df.to_string(index=False)
        docs.append(Document(
            page_content=conteudo,
            metadata={"source": caminho, "aba": nome_aba},
        ))
    return docs


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Uso: python loader.py <arquivo-ou-diretorio>")
        sys.exit(1)
    alvo = sys.argv[1]
    if os.path.isdir(alvo):
        docs = carregar_documentos_do_diretorio(alvo)
    else:
        docs = carregar_documento(alvo)
    print(f"\nTotal de documentos carregados: {len(docs)}")
    for d in docs[:3]:
        print(f"  -> {len(d.page_content)} chars | fonte: {d.metadata.get('source', '?')}")
